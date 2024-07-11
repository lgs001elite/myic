from pptx import Presentation
from pptx.util import Inches

# Load the presentation
file_path = 'figures_nsdi.pptx'
presentation = Presentation(file_path)

# Iterate through each slide
for slide in presentation.slides:
    # Determine the slide size
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height
    
    # Get the bounding box of all shapes on the slide
    min_x = min_y = float('inf')
    max_x = max_y = float('-inf')
    
    for shape in slide.shapes:
        if not shape.has_text_frame and not shape.has_table:
            continue
        min_x = min(min_x, shape.left)
        min_y = min(min_y, shape.top)
        max_x = max(max_x, shape.left + shape.width)
        max_y = max(max_y, shape.top + shape.height)
    
    # Calculate scale factors
    content_width = max_x - min_x
    content_height = max_y - min_y
    scale_x = slide_width / content_width
    scale_y = slide_height / content_height
    scale = min(scale_x, scale_y)
    
    # Center content
    offset_x = (slide_width - content_width * scale) / 2
    offset_y = (slide_height - content_height * scale) / 2
    
    # Scale and reposition shapes
    for shape in slide.shapes:
        shape.left = int((shape.left - min_x) * scale + offset_x)
        shape.top = int((shape.top - min_y) * scale + offset_y)
        shape.width = int(shape.width * scale)
        shape.height = int(shape.height * scale)

# Save the modified presentation
output_path = 'figures_nsdi_zoomed_in.pptx'
presentation.save(output_path)
